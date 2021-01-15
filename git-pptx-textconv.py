#!/usr/bin/env python3
# Original (Python 2): https://stackoverflow.com/a/32259984/14667736
# Prerequisite: $ pip install python-pptx

"""
Setup -- Add these lines to the following files:
--- .gitattributes
*.pptx diff=pptx

--- .gitconfig (or repo\.git\config    or your_user_home\.gitconfig) (change the path to point to your local copy of the script)
[diff "pptx"]
    binary = true
    textconv = python C:/Python27/Scripts/git-pptx-textconv.py

usage:
git diff your_powerpoint.pptx


Thanks to the  python-pptx docs and this snippet:
http://python-pptx.readthedocs.org/en/latest/user/quickstart.html#extract-all-text-from-slides-in-presentation
"""

import sys
from pptx import Presentation


if __name__ == '__main__':
    if len(sys.argv) != 2:
        print("Usage: git-pptx-textconv file.xslx")
        sys.exit(1)

    path_to_presentation = sys.argv[1]

    prs = Presentation(path_to_presentation)

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                par_text = ''
                for run in paragraph.runs:
                    s = run.text
                    s = s.replace(r"\\", "\\\\")
                    s = s.replace(r"\n", " ")
                    s = s.replace(r"\r", " ")
                    s = s.replace(r"\t", " ")
                    s = s.rstrip('\r\n')

                    # Convert left and right-hand quotes from Unicode to ASCII
                    # found http://stackoverflow.com/questions/816285/where-is-pythons-best-ascii-for-this-unicode-database
                    # go here if more power is needed  http://code.activestate.com/recipes/251871/
                    # or here                          https://pypi.python.org/pypi/Unidecode/0.04.1
                    punctuation = { 0x2018:0x27, 0x2019:0x27, 0x201C:0x22, 0x201D:0x22 }
                    s.translate(punctuation).encode('ascii', 'ignore')
                    if s:
                        par_text += s
                print(par_text)
